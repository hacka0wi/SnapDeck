# -*- coding: utf-8 -*-
"""Assemble ordered slide PNGs into a fresh 16:9 .pptx, each image full-bleed (edge to edge).

Usage:
    python3 assemble.py out.pptx order.json [slidesDir]

order.json = ["s01","s02",...]  (PNG basenames in slidesDir, without extension), in slide order.
If order.json is omitted, all slides/*.png are used sorted by name.
On macOS, if `import pptx` fails on python3, run with python3.11.
"""
import sys, os, json
from pptx import Presentation
from pptx.util import Inches

out = sys.argv[1] if len(sys.argv) > 1 else "out.pptx"
order_file = sys.argv[2] if len(sys.argv) > 2 else None
slides_dir = sys.argv[3] if len(sys.argv) > 3 else "slides"

if order_file and os.path.exists(order_file):
    order = json.load(open(order_file))
else:
    order = [f[:-4] for f in sorted(os.listdir(slides_dir)) if f.endswith(".png")]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
for key in order:
    png = os.path.join(slides_dir, f"{key}.png")
    if not os.path.exists(png):
        print("WARN missing", png); continue
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
prs.save(out)
print("saved", out, "·", len(prs.slides._sldIdLst), "slides")
